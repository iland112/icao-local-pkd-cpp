#!/usr/bin/env python3
"""
ICAO Local PKD 기술제안서 PPTX 생성 스크립트
- 템플릿: 전자여권위변조검사시스템(FASTpassPKD).pptx
- 디자인 참조: icao-local-pkd-proposal.html (21 slides)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os

# ── Constants ──
NAVY = RGBColor(0x02, 0x38, 0x5E)
NAVY_DEEP = RGBColor(0x01, 0x1D, 0x33)
NAVY_LIGHT = RGBColor(0x0A, 0x4F, 0x7A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x0F, 0x17, 0x2A)
GRAY = RGBColor(0x64, 0x74, 0x8B)
LIGHT_GRAY = RGBColor(0xE2, 0xE8, 0xF0)
GREEN = RGBColor(0x05, 0x96, 0x69)
RED = RGBColor(0xE1, 0x1D, 0x48)
BLUE = RGBColor(0x25, 0x63, 0xEB)
ORANGE = RGBColor(0xD9, 0x77, 0x06)
PURPLE = RGBColor(0x7C, 0x3A, 0xED)
BG_LIGHT = RGBColor(0xF8, 0xFA, 0xFC)
SURFACE = RGBColor(0xFF, 0xFF, 0xFF)
BORDER = RGBColor(0xE2, 0xE8, 0xF0)
ALT_ROW = RGBColor(0xEF, 0xF6, 0xFF)
BLUE_ACCENT = RGBColor(0x60, 0xA5, 0xFA)

SLIDE_W = Emu(10693400)
SLIDE_H = Emu(7561263)

FONT_KR = 'Malgun Gothic'

# ── Helper Functions ──

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _set_line_spacing(p, spacing):
    pPr = p._pPr
    if pPr is None:
        pPr = p._p.get_or_add_pPr()
    lnSpc = pPr.makeelement(qn('a:lnSpc'), {})
    spcPct = lnSpc.makeelement(qn('a:spcPct'), {'val': str(int(spacing * 100000))})
    lnSpc.append(spcPct)
    pPr.append(lnSpc)


def add_textbox(slide, left, top, width, height, text, font_size=12,
                font_color=BLACK, bold=False, alignment=PP_ALIGN.LEFT,
                font_name=FONT_KR, line_spacing=1.2, word_wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    _set_line_spacing(p, line_spacing)
    return txBox


def add_multiline_textbox(slide, left, top, width, height, lines,
                          font_size=11, font_color=BLACK, font_name=FONT_KR,
                          line_spacing=1.3, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, bold, color, size_override) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size_override or font_size)
        p.font.color.rgb = color or font_color
        p.font.bold = bold
        p.font.name = font_name
        p.alignment = alignment
        p.space_after = Pt(2)
    return txBox


def add_rounded_rect(slide, left, top, width, height, fill_color, text="",
                     font_size=10, font_color=WHITE, bold=False, border_color=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.bold = bold
        p.font.name = FONT_KR
    return shape


def add_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_table(slide, left, top, width, height, rows, cols):
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    return table_shape.table


def set_cell(table, row, col, text, font_size=9, font_color=BLACK,
             bold=False, fill_color=None, alignment=PP_ALIGN.LEFT):
    cell = table.cell(row, col)
    cell.text = text
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.bold = bold
        p.font.name = FONT_KR
        p.alignment = alignment
    if fill_color:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill_color
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE


def merge_cells(table, r1, c1, r2, c2):
    table.cell(r1, c1).merge(table.cell(r2, c2))


def add_page_header(slide, title, tag=None):
    add_rect(slide, Cm(0), Cm(0), SLIDE_W, Cm(0.15), NAVY)
    if tag:
        add_rounded_rect(slide, Cm(1.5), Cm(1.2), Cm(5.5), Cm(0.7), BLUE,
                         tag, font_size=8, bold=True)
    y = Cm(2.1) if tag else Cm(1.5)
    add_textbox(slide, Cm(1.5), y, Cm(24), Cm(1.2), title,
                font_size=22, font_color=BLACK, bold=True)
    add_rect(slide, Cm(1.5), y + Cm(1.2), Cm(4), Cm(0.08), NAVY)
    add_footer(slide)


def add_footer(slide):
    add_rect(slide, Cm(0), SLIDE_H - Cm(0.8), SLIDE_W, Cm(0.8), RGBColor(0xF1, 0xF5, 0xF9))
    add_textbox(slide, Cm(1.5), SLIDE_H - Cm(0.7), Cm(10), Cm(0.5),
                "FASTpass\u00ae PKD \u2014 \uc804\uc790\uc5ec\uad8c \uc704\xb7\ubcc0\uc870 \uac80\uc0ac \uc2dc\uc2a4\ud15c",
                font_size=7, font_color=GRAY)
    add_textbox(slide, Cm(20), SLIDE_H - Cm(0.7), Cm(6), Cm(0.5),
                "\u00a9 2026 SMARTCORE Inc.",
                font_size=7, font_color=GRAY, alignment=PP_ALIGN.RIGHT)


def add_stat_card(slide, left, top, width, height, value, label,
                  value_color=NAVY, bg_color=SURFACE):
    shape = add_rounded_rect(slide, left, top, width, height, bg_color,
                             border_color=BORDER)
    add_textbox(slide, left + Cm(0.3), top + Cm(0.3), width - Cm(0.6), Cm(1),
                value, font_size=20, font_color=value_color, bold=True,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left + Cm(0.3), top + Cm(1.5), width - Cm(0.6), Cm(0.7),
                label, font_size=8, font_color=GRAY,
                alignment=PP_ALIGN.CENTER)


def add_impact_card(slide, left, top, width, tag, desc, tag_color=BLUE, border_left_color=None):
    """Card with a small tag and description text, like the HTML impact-card."""
    shape = add_rounded_rect(slide, left, top, width, Cm(1.4), SURFACE,
                             border_color=BORDER)
    if border_left_color:
        # Add a thin left border accent
        add_rect(slide, left, top, Cm(0.1), Cm(1.4), border_left_color)
    # tag
    add_textbox(slide, left + Cm(0.4), top + Cm(0.15), Cm(3), Cm(0.5),
                tag, font_size=9, font_color=tag_color, bold=True)
    # desc
    add_textbox(slide, left + Cm(0.4), top + Cm(0.7), width - Cm(0.8), Cm(0.6),
                desc, font_size=8, font_color=GRAY, line_spacing=1.3)


def add_part_divider(prs, blank_layout, part_num, title, subtitle):
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, NAVY_DEEP)
    add_textbox(slide, Cm(3), Cm(6), Cm(20), Cm(1),
                f"Part {part_num}", font_size=14, font_color=RGBColor(0x94, 0xA3, 0xB8))
    add_rect(slide, Cm(3), Cm(7.2), Cm(6), Cm(0.08), BLUE_ACCENT)
    add_textbox(slide, Cm(3), Cm(7.5), Cm(20), Cm(2),
                title, font_size=32, font_color=WHITE, bold=True)
    add_textbox(slide, Cm(3), Cm(10), Cm(22), Cm(1),
                subtitle, font_size=12, font_color=RGBColor(0x64, 0x74, 0x8B))
    return slide


# ══════════════════════════════════════════════════════════════════════
# BUILD
# ══════════════════════════════════════════════════════════════════════

def build():
    template_path = os.path.join(os.path.dirname(__file__),
                                 "../references/\uc804\uc790\uc5ec\uad8c\uc704\ubcc0\uc870\uac80\uc0ac\uc2dc\uc2a4\ud15c(FASTpassPKD).pptx")
    prs = Presentation(template_path)

    # Remove all existing slides
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]

    blank_layout = prs.slide_layouts[6]  # blank

    # ================================================================
    # SLIDE 1: 표지 (Cover)
    # ================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, NAVY)
    add_rect(slide, Cm(18), Cm(0), Cm(10), Cm(3), NAVY_LIGHT)

    # Small label
    add_textbox(slide, Cm(3), Cm(3.5), Cm(10), Cm(0.6),
                "\uae30\uc220\uc81c\uc548\uc11c", font_size=9,
                font_color=RGBColor(0x94, 0xA3, 0xB8))

    # Badges
    badges_text = "ICAO Doc 9303  |  RFC 5280 X.509  |  BSI TR-03110"
    add_rounded_rect(slide, Cm(3), Cm(4.5), Cm(14), Cm(0.9),
                     RGBColor(0x0A, 0x4F, 0x7A),
                     badges_text, font_size=9, bold=False,
                     border_color=RGBColor(0x1A, 0x6F, 0xAA))

    # Title
    add_textbox(slide, Cm(3), Cm(6), Cm(22), Cm(2.5),
                "\uc804\uc790\uc5ec\uad8c \uc704\xb7\ubcc0\uc870 \uac80\uc0ac \uc2dc\uc2a4\ud15c",
                font_size=32, font_color=WHITE, bold=True, line_spacing=1.3)

    # Subtitle
    add_textbox(slide, Cm(3), Cm(8.5), Cm(22), Cm(1.5),
                "ICAO Local PKD \uad6c\uc131\xb7\uad00\ub9ac \ubc0f Passive Authentication \uac80\uc99d \ud1b5\ud569 \uc194\ub8e8\uc158",
                font_size=12, font_color=RGBColor(0x94, 0xA3, 0xB8), line_spacing=1.5)

    # Product name
    add_textbox(slide, Cm(3), Cm(10.5), Cm(20), Cm(1.5),
                "FASTpass\u00ae PKD",
                font_size=24, font_color=BLUE_ACCENT, bold=True)

    # Meta: 발주처 / 제안사 / 제출일
    add_textbox(slide, Cm(3), Cm(15), Cm(7), Cm(1.2),
                "\ubc1c\uc8fc\ucc98: \ubc95\ubb34\ubd80 \ucd9c\uc785\uad6d\xb7\uc678\uad6d\uc778\uc815\ucc45\ubcf8\ubd80\n\uc81c\uc548\uc0ac: (\uc8fc)\uc2a4\ub9c8\ud2b8\ucf54\uc5b4",
                font_size=10, font_color=RGBColor(0x94, 0xA3, 0xB8), line_spacing=1.6)
    add_textbox(slide, Cm(3), Cm(17), Cm(10), Cm(0.8),
                "2026\ub144 3\uc6d4", font_size=11,
                font_color=RGBColor(0x94, 0xA3, 0xB8))
    add_textbox(slide, Cm(17), Cm(17), Cm(8), Cm(0.8),
                "\u00a9 2026 SMARTCORE Inc.", font_size=10,
                font_color=RGBColor(0x64, 0x74, 0x8B), alignment=PP_ALIGN.RIGHT)

    # ================================================================
    # SLIDE 2: 목차
    # ================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_LIGHT)
    add_page_header(slide, "\ubaa9\ucc28")

    toc = [
        ("Part 1", "\uc0ac\uc5c5 \uc774\ud574 \ubc0f \ud604\ud669 \ubd84\uc11d",
         "\u2022 \uc81c\uc548 \ubc30\uacbd \ubc0f \ubaa9\uc801\n\u2022 \ud604\ud589 \uc2dc\uc2a4\ud15c \ud604\ud669\uacfc \ubb38\uc81c\uc810\n\u2022 \uc804\uc790\uc5ec\uad8c \ubcf4\uc548 \uc704\ud611 \ubd84\uc11d\n\u2022 PA \uac80\uc99d \ud544\uc694\uc131",
         RED, RGBColor(0xFE, 0xF2, 0xF2)),
        ("Part 2", "\uae30\uc220 \ubc29\uc548",
         "\u2022 \uc2dc\uc2a4\ud15c \uc544\ud0a4\ud14d\ucc98\n\u2022 \ud575\uc2ec \uae30\ub2a5 8\uac00\uc9c0\n\u2022 PA \uac80\uc99d \ud504\ub85c\uc138\uc2a4\n\u2022 \uc131\ub2a5 \ubca4\uce58\ub9c8\ud06c",
         BLUE, ALT_ROW),
        ("Part 3", "\ubcf4\uc548 \ubc0f \ud504\ub85c\uc81d\ud2b8 \uad00\ub9ac",
         "\u2022 \ubcf4\uc548 \ub300\ucc45 \ubc0f \ubc95\uc801 \uc900\uc218\n\u2022 \ucd94\uc9c4 \uccb4\uacc4 \ubc0f \uc77c\uc815\n\u2022 \uc720\uc9c0\ubcf4\uc218 \ubc0f \uae30\ub300\ud6a8\uacfc",
         GREEN, RGBColor(0xF0, 0xFD, 0xF4)),
    ]
    for i, (part, title, bullets, color, bg) in enumerate(toc):
        y = Cm(4.2 + i * 4)
        # Card background
        add_rounded_rect(slide, Cm(2), y, Cm(22), Cm(3.5), bg, border_color=BORDER)
        # Left border accent
        add_rect(slide, Cm(2), y, Cm(0.12), Cm(3.5), color)
        # Part label
        add_textbox(slide, Cm(2.8), y + Cm(0.3), Cm(4), Cm(0.5),
                    part, font_size=9, font_color=color, bold=True)
        # Title
        add_textbox(slide, Cm(2.8), y + Cm(0.9), Cm(10), Cm(0.7),
                    title, font_size=14, font_color=BLACK, bold=True)
        # Bullets
        add_textbox(slide, Cm(2.8), y + Cm(1.7), Cm(20), Cm(1.8),
                    bullets, font_size=9, font_color=GRAY, line_spacing=1.5)

    # ================================================================
    # SLIDE 3: Part 1 간지
    # ================================================================
    add_part_divider(prs, blank_layout, 1,
                     "\uc0ac\uc5c5 \uc774\ud574 \ubc0f \ud604\ud669 \ubd84\uc11d",
                     "\uc804\uc790\uc5ec\uad8c \ubcf4\uc548 \ud658\uacbd \ubcc0\ud654\uc640 PA \uac80\uc99d \ub3c4\uc785\uc758 \uc2dc\uae09\uc131")

    # ================================================================
    # SLIDE 4: 제안 배경 및 목적
    # ================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_LIGHT)
    add_page_header(slide, "\uc81c\uc548 \ubc30\uacbd \ubc0f \ubaa9\uc801", "Part 1 \u2014 \uc0ac\uc5c5 \uc774\ud574")

    # Left: 제안 배경
    add_textbox(slide, Cm(1.5), Cm(3.8), Cm(12), Cm(0.6),
                "\uc81c\uc548 \ubc30\uacbd", font_size=12, font_color=BLUE, bold=True)
    bg_items = [
        ("* \uc2dc\uc2a4\ud15c \ub178\ud6c4\ud654", "\ud604 PKD Manager 2009\ub144 \ub3c4\uc785, HW/SW \ub0b4\uc6a9\uc5f0\uc218 \ucd08\uacfc \xb7 \uc9c0\uc6d0 \ub9cc\ub8cc"),
        ("* \uae30\uc220 \ud45c\uc900 \ubcc0\uacbd", "LDS 1.8, PACE \ub4f1 \uc2e0\uaddc \uae30\uc220 \ubbf8\uc9c0\uc6d0 (ICAO 2028\ub144 \uc758\ubb34 \uc801\uc6a9 \uc608\uc815)"),
        ("* \uc704\uc870 \uae30\uc220 \uace0\ub3c4\ud654", "BAC \ud1b5\uacfc + \uce69 \ub370\uc774\ud130 \uc815\uc0c1 \uc77d\uae30\uae4c\uc9c0 \uac00\ub2a5\ud55c 3\uc138\ub300 \uc815\ubc00 \uc704\uc870 \uae09\uc99d"),
        ("* \uc11c\ube44\uc2a4 \ubc94\uc704 \ud655\ub300", "\uc785\xb7\ucd9c\uad6d\uc2ec\uc0ac \uc678 \uc778\ucc9c\uacf5\ud56d\uacf5\uc0ac, \uc7ac\uc678\uacf5\uad00 \ub4f1 \ud655\uc7a5 \ud544\uc694"),
    ]
    for i, (title, desc) in enumerate(bg_items):
        y = Cm(4.6 + i * 1.5)
        add_textbox(slide, Cm(1.5), y, Cm(12), Cm(0.5),
                    title, font_size=9, font_color=BLACK, bold=True)
        add_textbox(slide, Cm(1.5), y + Cm(0.5), Cm(12), Cm(0.7),
                    desc, font_size=8, font_color=GRAY)

    # Right: 사업 목적
    add_textbox(slide, Cm(14), Cm(3.8), Cm(12), Cm(0.6),
                "\uc0ac\uc5c5 \ubaa9\uc801", font_size=12, font_color=GREEN, bold=True)
    purposes = [
        ("1. \uc804\uc790\uc5ec\uad8c\uc2dc\uc2a4\ud15c \uc804\uba74 \uc7ac\uad6c\ucd95",
         "\ubcc0\uacbd\ub41c \uae30\uc220 \ud45c\uc900(LDS 1.8, PACE) \ubc18\uc601 \ubc0f \uc778\uc99d \uccb4\uacc4 \uace0\ub3c4\ud654"),
        ("2. Local PKD (\uc7ac)\uad6c\ucd95",
         "ICAO PKD\u2013Local PKD \uc5f0\uacc4, PKD \uccb4\uacc4 \ud65c\uc6a9\uc131 \uac15\ud654"),
        ("3. \uc804\uc790\uce69 \ubd84\uc11d \ub2a5\ub825 \uac15\ud654",
         "\uc815\ubc00\ubd84\uc11d\xb7\ud328\ud134\ubd84\uc11d\xb7\uacb0\ud568\ucc98\ub9ac \ud504\ub85c\uadf8\ub7a8\uc73c\ub85c \uac80\uc99d \uacb0\uacfc \uc2e0\ub8b0\ub3c4 \ud5a5\uc0c1"),
        ("4. \uc6b4\uc601 \uad00\ub9ac\uc131\xb7\uac00\uc2dc\uc131 \ud655\ubcf4",
         "\ubaa8\ub2c8\ud130\ub9c1 \uc5d0\uc774\uc804\ud2b8, \uc6b4\uc601 \ud1b5\uacc4, \ubc1c\uae09\uc0c1 \uacb0\ud568\uc815\ubcf4 \uad00\ub9ac \uccb4\uacc4"),
    ]
    for i, (title, desc) in enumerate(purposes):
        y = Cm(4.6 + i * 1.5)
        add_rect(slide, Cm(14), y, Cm(0.1), Cm(1.2), GREEN)
        add_textbox(slide, Cm(14.5), y, Cm(11), Cm(0.5),
                    title, font_size=9, font_color=BLACK, bold=True)
        add_textbox(slide, Cm(14.5), y + Cm(0.5), Cm(11), Cm(0.7),
                    desc, font_size=8, font_color=GRAY)

    # ================================================================
    # SLIDE 5: 현행 시스템 현황과 문제점
    # ================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_LIGHT)
    add_page_header(slide, "\ud604\ud589 \uc804\uc790\uc5ec\uad8c\uc2dc\uc2a4\ud15c \ud604\ud669\uacfc \ubb38\uc81c\uc810", "Part 1 \u2014 \ud604\ud669 \ubd84\uc11d")

    # Left: 현행 운영 현황
    add_textbox(slide, Cm(1.5), Cm(3.8), Cm(10), Cm(0.6),
                "\ud604\ud589 \uc6b4\uc601 \ud604\ud669", font_size=11, font_color=BLUE, bold=True)
    current_items = [
        "\u2713 ICRM \uae30\ubc18 \uc785\xb7\ucd9c\uad6d\uc2ec\uc0ac \uc2dc \uc804\uc790\uc5ec\uad8c \ud310\ub3c5\xb7\uc778\uc99d",
        "\u2713 BAC/CA/AA \uc778\uc99d + PKI Toolkit PA \uac80\uc99d",
        "\u2713 \uc790\ub3d9\ucd9c\uc785\uad6d\uc2ec\uc0ac(SES) 18\uac1c\uad6d \uc6b4\uc601",
        "\u2713 PKD Manager\ub85c ICAO PKD \u2192 Local PKD \ub3d9\uae30",
    ]
    for i, item in enumerate(current_items):
        add_textbox(slide, Cm(1.5), Cm(4.5 + i * 0.7), Cm(12), Cm(0.6),
                    item, font_size=9, font_color=BLACK)

    # Left: 핵심 문제점
    add_textbox(slide, Cm(1.5), Cm(7.5), Cm(10), Cm(0.6),
                "\ud575\uc2ec \ubb38\uc81c\uc810", font_size=11, font_color=RED, bold=True)
    problems = [
        "\u2717 \uc2dc\uc2a4\ud15c \ub178\ud6c4\ud654 \u2014 PKD Manager 2009\ub144 \ub3c4\uc785, HW/SW 6\ub144+ \ub0b4\uc6a9\uc5f0\uc218 \ucd08\uacfc",
        "\u2717 \uae30\uc220 \ud45c\uc900 \ubbf8\uc9c0\uc6d0 \u2014 LDS 1.8, PACE \ubbf8\uc9c0\uc6d0 (ICAO 2028\ub144 \uc758\ubb34 \uc801\uc6a9 \uad8c\uace0)",
        "\u2717 ML \ubbf8\ud65c\uc6a9 \u2014 ICAO Master List \ud1b5\ud55c CSCA \uc778\uc99d\uc11c \ubcf4\uc870 \uc218\uc9d1\xb7\uac80\uc99d \ubd80\uc7ac",
        "\u2717 DL\xb7NC \ubbf8\ud65c\uc6a9 \u2014 \uac01\uad6d DL \ubc0f \ube44\ud45c\uc900 DSC \ubbf8\uc900\uc218 \uc0ac\ud56d \uad00\ub9ac\xb7\uc870\ud68c \ubd80\uc7ac",
    ]
    for i, item in enumerate(problems):
        add_textbox(slide, Cm(1.5), Cm(8.2 + i * 0.8), Cm(12), Cm(0.6),
                    item, font_size=8, font_color=RGBColor(0x4B, 0x55, 0x63))

    # Right: 추가 개선 필요 사항
    add_textbox(slide, Cm(14), Cm(3.8), Cm(12), Cm(0.6),
                "\ucd94\uac00 \uac1c\uc120 \ud544\uc694 \uc0ac\ud56d", font_size=11, font_color=ORANGE, bold=True)
    improvements = [
        ("\uac80\uc99d", "\uac80\uc99d \uacb0\uacfc \uc778\uc2dd\uc131 \ubd80\uc871 \u2014 \uc131\uacf5/\uc2e4\ud328 \uc774\ubd84\ubc95, \uc2e0\ub8b0\ub3c4 \uac1c\ub150 \ubd80\uc7ac"),
        ("\ubd84\uc11d", "\uc804\uc790\uce69 \ubd84\uc11d\ub3c4\uad6c \ubd80\uc7ac \u2014 \uc704\xb7\ubcc0\uc870 \uc758\uc2ec \uc2dc \uc0c1\uc138 \ubd84\uc11d \uc18c\ud504\ud2b8\uc6e8\uc5b4 \uc5c6\uc74c"),
        ("\ubaa8\ub2c8\ud130\ub9c1", "\ubaa8\ub2c8\ud130\ub9c1 \uccb4\uacc4 \ubd80\uc7ac \u2014 \uc624\ub958 \ube48\ub3c4\xb7\uc720\ud615 \uc0ac\uc804 \ud0d0\uc9c0 \uae30\ub2a5 \ubbf8\uad6c\ud604"),
        ("\uad00\ub9ac", "\uc778\uc99d\uc11c \uad00\ub9ac \uc81c\uc57d \u2014 PKD Manager \ubcc4\ub3c4 \uad00\ub9ac \uae30\ub2a5 \ubbf8\uc81c\uacf5"),
        ("\ubc94\uc704", "\uc11c\ube44\uc2a4 \ubc94\uc704 \ud55c\uc815 \u2014 \ucd9c\uc785\uad6d\uc2ec\uc0ac\ub9cc \ub300\uc0c1, \uc778\ucc9c\uacf5\ud56d\uacf5\uc0ac \ub4f1 \ud655\uc7a5 \ud544\uc694"),
    ]
    for i, (tag, desc) in enumerate(improvements):
        y = Cm(4.6 + i * 1.5)
        add_rounded_rect(slide, Cm(14), y, Cm(2.2), Cm(0.7), ALT_ROW,
                         tag, font_size=8, font_color=BLUE, bold=True,
                         border_color=BORDER)
        add_textbox(slide, Cm(16.5), y + Cm(0.05), Cm(9), Cm(0.7),
                    desc, font_size=8, font_color=BLACK)

    # ================================================================
    # SLIDE 6: 사업 과제 대응 범위
    # ================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_LIGHT)
    add_page_header(slide, "\uc0ac\uc5c5 \uacfc\uc81c \ub300\uc751 \ubc94\uc704", "Part 1 \u2014 \uc0ac\uc5c5 \uc774\ud574")

    add_textbox(slide, Cm(1.5), Cm(3.5), Cm(24), Cm(0.5),
                "\uc804\uc790\uc5ec\uad8c\uc2dc\uc2a4\ud15c \uc7ac\uad6c\ucd95 \uc0ac\uc5c5\uacc4\ud68d\uc11c 22\uac1c \uacfc\uc81c \uc911 FASTpass\u00ae PKD \uc9c1\uc811 \ub300\uc751 \ud56d\ubaa9",
                font_size=9, font_color=GRAY)

    task_rows = [
        ("\uc804\uc790\uc5ec\uad8c \uc778\uc99d(PA) \uacf5\ud1b5 \ud504\ub85c\uadf8\ub7a8", "PA Service \u2014 ICAO 9303 8\ub2e8\uacc4 \uac80\uc99d", "\uc644\uc804 \ub300\uc751"),
        ("Local PKD \uad6c\ucd95", "ICAO PKD + \uc591\uc790\ud611\uc815 + DB/LDAP \ud1b5\ud569 \uc800\uc7a5\uc18c", "\uc644\uc804 \ub300\uc751"),
        ("Local PKD \uad00\ub9ac \ud504\ub85c\uadf8\ub7a8", "PKD Management \u2014 ML/CSCA/DS/CRL/DL \uad00\ub9ac", "\uc644\uc804 \ub300\uc751"),
        ("\uc804\uc790\uc5ec\uad8c \uc6b4\uc601\ud658\uacbd \uad00\ub9ac", "\ud658\uacbd\uc124\uc815, \uc778\uc99d \uc801\uc6a9 \uc5ec\ubd80 \uad00\ub9ac", "\ubd80\ubd84 \ub300\uc751"),
        ("\uc804\uc790\uc5ec\uad8c \uc815\ubc00\ubd84\uc11d \ud504\ub85c\uadf8\ub7a8", "\uc778\uc99d\uc11c \uc9c1\uc811 \ud30c\uc2f1\xb7\ubd84\uc11d, Doc 9303 \uc900\uc218 \uac80\uc0ac", "\uc644\uc804 \ub300\uc751"),
        ("\uc804\uc790\uc5ec\uad8c \uc778\uc99d \uacb0\ud568\ucc98\ub9ac", "Doc 9303 Compliance Checklist (28\uac1c \ud56d\ubaa9)", "\uc644\uc804 \ub300\uc751"),
        ("\ub370\uc774\ud130 \uc218\uc9d1 \ubc0f \ud328\ud134\ubd84\uc11d", "\uac80\uc99d \uacb0\uacfc \ub370\uc774\ud130 \uc218\uc9d1 \xb7 \uad6d\uac00\ubcc4/\uc720\ud615\ubcc4 \ud1b5\uacc4 \ubd84\uc11d", "\uc644\uc804 \ub300\uc751"),
        ("\ube44\uc815\uc0c1 \ud328\ud134 \ub370\uc774\ud130 \uad00\ub9ac", "\uc624\ub958 \uc720\ud615\ubcc4 \ubaa9\ub85d \uc870\ud68c \xb7 \ube48\ub3c4/\uadf8\ub798\ud504 \ud1b5\uacc4", "\uc644\uc804 \ub300\uc751"),
        ("\ud310\ub3c5\xb7\uc778\uc99d \uc624\ub958 \ub370\uc774\ud130 \uad00\ub9ac", "\uc624\ub958 \uac74\ubcc4 \uc0c1\uc138\uc870\ud68c \xb7 \ubc14\uc774\ub108\ub9ac \ub370\uc774\ud130 \ub2e4\uc6b4\ub85c\ub4dc", "\uc644\uc804 \ub300\uc751"),
        ("\ubc1c\uae09\uc0c1 \uacb0\ud568\uc815\ubcf4 \uad00\ub9ac", "DL/Non-Conformant \uc870\ud68c\xb7\uad00\ub9ac \uae30\ub2a5", "\uc644\uc804 \ub300\uc751"),
        ("\ud310\ub3c5\xb7\uc778\uc99d \uc6b4\uc601 \ud1b5\uacc4 \uc870\ud68c", "Dashboard + \ubcf4\uace0\uc11c (DSC_NC/CRL/Trust Chain)", "\uc644\uc804 \ub300\uc751"),
        ("\ud310\ub3c5\xb7\uc778\uc99d \ubaa8\ub2c8\ud130\ub9c1 \uc5d0\uc774\uc804\ud2b8", "Monitoring Service + SSE \uc2e4\uc2dc\uac04 \uc54c\ub9bc", "\uc644\uc804 \ub300\uc751"),
        ("\uc785\xb7\ucd9c\uad6d\uc2ec\uc0ac \uc2dc \ud310\ub3c5\xb7\uc778\uc99d", "REST API + X-API-Key \uc778\uc99d \uc5f0\ub3d9", "API \uc81c\uacf5"),
    ]
    num_rows = len(task_rows) + 1
    tbl = add_table(slide, Cm(1.5), Cm(4.1), Cm(24), Cm(8.5), num_rows, 3)
    tbl.columns[0].width = Cm(8.5)
    tbl.columns[1].width = Cm(10.5)
    tbl.columns[2].width = Cm(5)
    set_cell(tbl, 0, 0, "\uc0ac\uc5c5 \uacfc\uc81c", bold=True, fill_color=NAVY, font_color=WHITE, font_size=8, alignment=PP_ALIGN.CENTER)
    set_cell(tbl, 0, 1, "FASTpass\u00ae PKD \ub300\uc751 \uae30\ub2a5", bold=True, fill_color=NAVY, font_color=WHITE, font_size=8, alignment=PP_ALIGN.CENTER)
    set_cell(tbl, 0, 2, "\uc0c1\ud0dc", bold=True, fill_color=NAVY, font_color=WHITE, font_size=8, alignment=PP_ALIGN.CENTER)
    for i, (task, func, status) in enumerate(task_rows):
        r = i + 1
        bg = ALT_ROW if r % 2 == 1 else None
        s_color = GREEN if "\uc644\uc804" in status else BLUE
        set_cell(tbl, r, 0, task, font_size=7, bold=True, fill_color=bg)
        set_cell(tbl, r, 1, func, font_size=7, fill_color=bg)
        set_cell(tbl, r, 2, status, font_size=7, font_color=s_color, bold=True, fill_color=bg, alignment=PP_ALIGN.CENTER)

    # Bottom stats
    add_stat_card(slide, Cm(7), Cm(13), Cm(4.5), Cm(2), "11/22", "\uc644\uc804 \ub300\uc751 \uacfc\uc81c", GREEN)
    add_stat_card(slide, Cm(14), Cm(13), Cm(4.5), Cm(2), "13/22", "\uc9c1\uc811 \ub300\uc751 \uacfc\uc81c", BLUE)

    # ================================================================
    # SLIDE 7: 위·변조 위협 분석
    # ================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_LIGHT)
    add_page_header(slide, "\uc804\uc790\uc5ec\uad8c \uc704\xb7\ubcc0\uc870 \uc704\ud611 \ubd84\uc11d", "Part 1 \u2014 \ud604\ud669 \ubd84\uc11d")

    stats = [
        ("3,200", "2021\ub144 \uc801\ubc1c", BLACK),
        ("4,500", "2022\ub144 (\u219141%)", ORANGE),
        ("5,800", "2023\ub144 (\u219129%)", RED),
        ("7,000+", "2024\ub144 \ucd94\uc815", RED),
    ]
    for i, (val, label, color) in enumerate(stats):
        add_stat_card(slide, Cm(1.5 + i * 6.2), Cm(4), Cm(5.7), Cm(2.5), val, label, color)

    add_textbox(slide, Cm(1.5), Cm(7.5), Cm(20), Cm(0.6),
                "\uc704\uc870 \uae30\uc220 3\uc138\ub300 \uc9c4\ud654", font_size=13, font_color=BLACK, bold=True)

    gens = [
        ("1\uc138\ub300", "MRZ \uc704\uc870", "\uae30\uacc4 \ud310\ub3c5 \uc601\uc5ed \uc870\uc791 (\uc721\uc548 \uac80\uc0ac\ub85c \ubc1c\uacac)", ORANGE),
        ("2\uc138\ub300", "\uce69 \ubcf5\uc81c", "\uc6d0\ubcf8 \ub370\uc774\ud130 \ubcf5\uc0ac (BAC \ud1b5\uacfc)", RED),
        ("3\uc138\ub300", "\uc815\ubc00 \uc704\uc870", "BAC\xb7\uc721\uc548 \ubaa8\ub450 \ud1b5\uacfc, PA\ub9cc\uc774 \ud0d0\uc9c0 \uac00\ub2a5", PURPLE),
    ]
    for i, (gen, title, desc, color) in enumerate(gens):
        y = Cm(8.5 + i * 2)
        add_rounded_rect(slide, Cm(1.5), y, Cm(2.5), Cm(1), color,
                         gen, font_size=9, bold=True)
        add_textbox(slide, Cm(4.5), y + Cm(0.05), Cm(20), Cm(1),
                    f"{title} \u2014 {desc}", font_size=10, font_color=BLACK)

    # ================================================================
    # SLIDE 8: PA 검증 필요성
    # ================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_LIGHT)
    add_page_header(slide, "Passive Authentication \uac80\uc99d \uccb4\uacc4 \ub3c4\uc785\uc758 \ud544\uc694\uc131", "Part 1 \u2014 \uac1c\uc120 \ubc29\ud5a5")

    # Left: ICAO quotes + PA guarantees
    add_rounded_rect(slide, Cm(1.5), Cm(4), Cm(12), Cm(1.5), SURFACE, border_color=BORDER)
    add_rect(slide, Cm(1.5), Cm(4), Cm(0.1), Cm(1.5), ORANGE)
    add_textbox(slide, Cm(2), Cm(4.1), Cm(11), Cm(0.4),
                "2010\ub144\ub300 \u2014 ICAO \uad8c\uace0", font_size=7, font_color=GRAY)
    add_textbox(slide, Cm(2), Cm(4.5), Cm(11), Cm(0.6),
                '"Passive Authentication is recommended"', font_size=10, font_color=BLACK, bold=False)

    add_rounded_rect(slide, Cm(1.5), Cm(5.8), Cm(12), Cm(1.5), SURFACE, border_color=BORDER)
    add_rect(slide, Cm(1.5), Cm(5.8), Cm(0.1), Cm(1.5), GREEN)
    add_textbox(slide, Cm(2), Cm(5.9), Cm(11), Cm(0.4),
                "2024\ub144 \u2014 ICAO Doc 9303 \uac1c\uc815", font_size=7, font_color=GRAY)
    add_textbox(slide, Cm(2), Cm(6.3), Cm(11), Cm(0.6),
                '"Passive Authentication MUST be implemented"', font_size=10, font_color=GREEN, bold=True)

    add_textbox(slide, Cm(1.5), Cm(7.8), Cm(12), Cm(0.5),
                "PA \uac80\uc99d\uc774 \ubcf4\uc7a5\ud558\ub294 \uac83", font_size=10, font_color=BLACK, bold=True)
    pa_items = [
        "\u2713 \uce69 \ub370\uc774\ud130 \ubb34\uacb0\uc131 \u2014 SOD \uc11c\uba85\uc73c\ub85c \ubcc0\uc870 \ud0d0\uc9c0",
        "\u2713 \ubc1c\uae09\uad6d \uc778\uc99d \u2014 CSCA \u2192 DSC Trust Chain \uac80\uc99d",
        "\u2713 \uc778\uc99d\uc11c \uc720\ud6a8\uc131 \u2014 CRL \uae30\ubc18 \ud3d0\uae30 \uc5ec\ubd80 \ud655\uc778",
        "\u2713 DG \ud574\uc2dc \ubb34\uacb0\uc131 \u2014 \uc0ac\uc9c4(DG2)\xb7MRZ(DG1) \ub4f1 \uc6d0\ubcf8 \ubcf4\uc7a5",
    ]
    for i, item in enumerate(pa_items):
        add_textbox(slide, Cm(1.5), Cm(8.5 + i * 0.7), Cm(12), Cm(0.5),
                    item, font_size=8, font_color=BLACK)

    # Right: PA flow + stats
    add_textbox(slide, Cm(14), Cm(4), Cm(12), Cm(0.5),
                "PA \uac80\uc99d \ud504\ub85c\uc138\uc2a4", font_size=10, font_color=BLACK, bold=True)
    steps = [
        ("Step 1-2", "SOD \uc11c\uba85 \uac80\uc99d + DSC \uc870\ud68c", BLUE),
        ("Step 3-4", "CSCA \uc870\ud68c + Trust Chain \uac80\uc99d", BLUE),
        ("Step 5-7", "SOD \ud574\uc2dc \ucd94\ucd9c + DG \ud574\uc2dc \ube44\uad50", ORANGE),
        ("Step 8", "CRL \ud3d0\uae30 \uc5ec\ubd80 \ud655\uc778", GREEN),
    ]
    for i, (step, desc, color) in enumerate(steps):
        y = Cm(4.8 + i * 1.3)
        add_rounded_rect(slide, Cm(14), y, Cm(2.5), Cm(0.8), color,
                         step, font_size=7, bold=True)
        add_textbox(slide, Cm(17), y + Cm(0.1), Cm(8.5), Cm(0.6),
                    desc, font_size=9, font_color=BLACK)

    # 3 stat cards
    add_stat_card(slide, Cm(14), Cm(10.5), Cm(3.5), Cm(2), "99.9%", "\uac80\uc99d \uc815\ud655\ub3c4", GREEN)
    add_stat_card(slide, Cm(18), Cm(10.5), Cm(3.5), Cm(2), "100-300ms", "\ucc98\ub9ac \uc2dc\uac04", NAVY)
    add_stat_card(slide, Cm(22), Cm(10.5), Cm(3.5), Cm(2), "8\ub2e8\uacc4", "\uc790\ub3d9 \uac80\uc99d", BLUE)

    # ================================================================
    # SLIDE 9: Part 2 간지
    # ================================================================
    slide = add_part_divider(prs, blank_layout, 2,
                     "\uae30\uc220 \ubc29\uc548",
                     "\ub9c8\uc774\ud06c\ub85c\uc11c\ube44\uc2a4 \uc544\ud0a4\ud14d\ucc98 \xb7 ICAO Doc 9303 \ub4f1 \uad00\ub828 \uad6d\uc81c \ud45c\uc900 \uc900\uc218 \xb7 \ucee8\ud14c\uc774\ub108 \uae30\ubc18 \ubc30\ud3ec \xb7 HA \uc774\uc911\ud654 \uad6c\uc131")
    # Add FASTpass® PKD title
    add_textbox(slide, Cm(3), Cm(9.5), Cm(20), Cm(1.5),
                "FASTpass\u00ae PKD",
                font_size=22, font_color=BLUE_ACCENT, bold=True)

    # ================================================================
    # SLIDE 10: 시스템 아키텍처
    # ================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_LIGHT)
    add_page_header(slide, "\uc2dc\uc2a4\ud15c \uc544\ud0a4\ud14d\ucc98", "Part 2 \u2014 \uc2dc\uc2a4\ud15c \uc124\uacc4")

    # Left: 4 microservices
    add_textbox(slide, Cm(1.5), Cm(3.8), Cm(12), Cm(0.6),
                "4\uac1c \ub3c5\ub9bd \ub9c8\uc774\ud06c\ub85c\uc11c\ube44\uc2a4", font_size=12, font_color=BLACK, bold=True)
    services = [
        ("PKD", "PKD Management", "\uc778\uc99d\uc11c \uc5c5\ub85c\ub4dc\xb7\ud30c\uc2f1\xb7\uac80\uc99d\xb7\uad00\ub9ac", BLUE),
        ("PA", "PA Service", "ICAO 9303 8\ub2e8\uacc4 \uc704\xb7\ubcc0\uc870 \uac80\uc99d", ORANGE),
        ("Sync", "PKD Relay", "DB\u2194LDAP \uc790\ub3d9 \ub3d9\uae30\ud654\xb7\uc870\uc815", GREEN),
        ("Mon", "Monitoring", "\uc2dc\uc2a4\ud15c \uba54\ud2b8\ub9ad\xb7\uc11c\ube44\uc2a4 \ud5ec\uc2a4\uccb4\ud06c", PURPLE),
    ]
    for i, (abbr, name, desc, color) in enumerate(services):
        y = Cm(4.8 + i * 1.8)
        add_rounded_rect(slide, Cm(1.5), y, Cm(2), Cm(0.9), color,
                         abbr, font_size=8, bold=True)
        add_textbox(slide, Cm(4), y, Cm(9), Cm(0.5),
                    name, font_size=10, font_color=BLACK, bold=True)
        add_textbox(slide, Cm(4), y + Cm(0.5), Cm(9), Cm(0.5),
                    desc, font_size=9, font_color=GRAY)

    # Right: tech stack
    add_textbox(slide, Cm(14), Cm(3.8), Cm(12), Cm(0.6),
                "\uae30\uc220 \uc2a4\ud0dd", font_size=12, font_color=BLACK, bold=True)
    tech_rows = [
        ("Backend", "C++20 / Drogon Framework"),
        ("Frontend", "React 19 / TypeScript / Tailwind CSS"),
        ("Database", "Oracle RDBMS"),
        ("Directory", "OpenLDAP (MMR \uc774\uc911\ud654)"),
        ("Gateway", "nginx (SSL/TLS, Rate Limiting)"),
        ("Container", "Docker / Podman (RHEL 9)"),
        ("\uc554\ud638 \ub77c\uc774\ube0c\ub7ec\ub9ac", "OpenSSL 3.x (FIPS \uc9c0\uc6d0)"),
    ]
    tbl = add_table(slide, Cm(14), Cm(4.8), Cm(11.5), Cm(5.6), len(tech_rows) + 1, 2)
    tbl.columns[0].width = Cm(4)
    tbl.columns[1].width = Cm(7.5)
    set_cell(tbl, 0, 0, "\uad6c\ubd84", bold=True, fill_color=NAVY, font_color=WHITE, font_size=9, alignment=PP_ALIGN.CENTER)
    set_cell(tbl, 0, 1, "\uae30\uc220", bold=True, fill_color=NAVY, font_color=WHITE, font_size=9, alignment=PP_ALIGN.CENTER)
    for i, (k, v) in enumerate(tech_rows):
        set_cell(tbl, i + 1, 0, k, bold=True, font_size=9, alignment=PP_ALIGN.CENTER, fill_color=ALT_ROW)
        set_cell(tbl, i + 1, 1, v, font_size=9)

    # ================================================================
    # SLIDE 11: SW 아키텍처 구성도
    # ================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_LIGHT)
    add_page_header(slide, "\uc18c\ud504\ud2b8\uc6e8\uc5b4 \uc544\ud0a4\ud14d\ucc98 \uad6c\uc131\ub3c4", "Part 2 \u2014 \uc2dc\uc2a4\ud15c \uc124\uacc4")

    # 3-tier architecture as table
    arch_rows = [
        ("\ud504\ub860\ud2b8\uc5d4\ub4dc", "React 19 / TypeScript / Tailwind CSS / Vite (27\uac1c \ud398\uc774\uc9c0)"),
        ("API Gateway", "nginx \u2014 SSL/TLS \uc885\ub2e8, Rate Limiting, \ub85c\ub4dc \ubc38\ub7f0\uc2f1, CSP \ud5e4\ub354, auth_request"),
        ("PKD Management :8081", "\uc778\uc99d\uc11c \uc5c5\ub85c\ub4dc\xb7\ud30c\uc2f1 / Trust Chain \uac80\uc99d / ICAO 9303 \uc900\uc218 \uac80\uc0ac / CSR \uad00\ub9ac / API Client \uc778\uc99d"),
        ("PA Service :8082", "SOD \uc11c\uba85 \uac80\uc99d / DSC\u2192CSCA \uccb4\uc778 / DG \ud574\uc2dc \ube44\uad50 / CRL \ud3d0\uae30 \ud655\uc778 / DSC \uc790\ub3d9 \ub4f1\ub85d"),
        ("PKD Relay :8083", "DB\u2194LDAP \ub3d9\uae30\ud654 / \uc790\ub3d9 Reconciliation / DSC \uc7ac\uac80\uc99d / SSE \uc2e4\uc2dc\uac04 \uc54c\ub9bc"),
        ("Monitoring :8084", "\uc11c\ube44\uc2a4 \ud5ec\uc2a4\uccb4\ud06c / \uc2dc\uc2a4\ud15c \uba54\ud2b8\ub9ad / \uc811\uc18d\uc790 \ud604\ud669 / \uc2e4\uc2dc\uac04 \uc54c\ub9bc"),
        ("Oracle RDBMS", "IQueryExecutor \ud328\ud134 \xb7 Repository Pattern \xb7 certificate \xb7 validation_result \xb7 audit_log"),
        ("LDAP (OpenLDAP MMR)", "Primary + Replica \uc774\uc911\ud654 \xb7 dc=data (CSCA/DSC/CRL/ML) \xb7 dc=nc-data (DSC_NC)"),
        ("\uc678\ubd80 \uc5f0\uacc4", "ICRM \uc2ec\uc0ac PC (1,500\ub300) / \uc5f0\uacc4\uae30\uad00 (REST API) / ICAO PKD (CSR+LDAP/REST) / \uc678\uad50\ubd80"),
    ]
    tbl = add_table(slide, Cm(1.5), Cm(3.8), Cm(24), Cm(9), len(arch_rows) + 1, 2)
    tbl.columns[0].width = Cm(6)
    tbl.columns[1].width = Cm(18)
    set_cell(tbl, 0, 0, "\uacc4\uce35", bold=True, fill_color=NAVY, font_color=WHITE, font_size=9, alignment=PP_ALIGN.CENTER)
    set_cell(tbl, 0, 1, "\uad6c\uc131 \uc694\uc18c", bold=True, fill_color=NAVY, font_color=WHITE, font_size=9, alignment=PP_ALIGN.CENTER)
    for i, (layer, detail) in enumerate(arch_rows):
        bg = ALT_ROW if i % 2 == 0 else None
        set_cell(tbl, i + 1, 0, layer, bold=True, font_size=8, fill_color=bg, alignment=PP_ALIGN.CENTER)
        set_cell(tbl, i + 1, 1, detail, font_size=8, fill_color=bg)

    # Design patterns badges
    patterns = "Repository Pattern \xb7 Query Executor (IQueryExecutor) \xb7 Service Container (DI) \xb7 Handler Pattern \xb7 RAII \xb7 Factory Pattern \xb7 SSE \xb7 Docker/Podman"
    add_textbox(slide, Cm(1.5), Cm(13.2), Cm(24), Cm(0.5),
                patterns, font_size=7, font_color=GRAY, alignment=PP_ALIGN.CENTER)

    # ================================================================
    # SLIDE 12: H/W 네트워크 구성도
    # ================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_LIGHT)
    add_page_header(slide, "H/W \ub124\ud2b8\uc6cc\ud06c \uad6c\uc131\ub3c4 (HA \uc774\uc911\ud654)", "Part 2 \u2014 \uc2dc\uc2a4\ud15c \uc124\uacc4")

    zones = [
        ("\ub0b4\ubd80\ub9dd", [
            "\uad00\ub9ac\uc11c\ubc84 x1 \u2014 API Gateway, \ubaa8\ub2c8\ud130\ub9c1, UI (16Core/32GB)",
            "PKD\uc11c\ubc84 #1 Active \u2014 PA, PKD \uad00\ub9ac, LDAP Primary (16Core/32GB)",
            "PKD\uc11c\ubc84 #2 Active \u2014 PA, PKD \uad00\ub9ac, LDAP Replica (16Core/32GB)",
            "SAN \uc2a4\uc704\uce58 #1/#2 \u2014 32Gb FC, Fabric A/B \uc774\uc911\ud654",
            "\uc2a4\ud1a0\ub9ac\uc9c0 \u2014 NVMe 7.68TB\xd728EA, 150.42TiB, RaidB",
            "Oracle RDBMS + LDAP \uc774\uc911\ud654",
        ], BLUE),
        ("DMZ", [
            "\uc911\uacc4\uc11c\ubc84 #1 Active \u2014 ICAO PKD \uc5f0\uacc4 (CSR+LDAP/REST)",
            "\uc911\uacc4\uc11c\ubc84 #2 Standby \u2014 ICAO PKD \uc5f0\uacc4",
            "Active-Standby \uc774\uc911\ud654",
        ], ORANGE),
        ("\uc678\ubd80", [
            "ICAO PKD \u2014 CSCA/DSC/CRL/DL/NC-DATA (CSR+LDAP/REST)",
            "\uc678\uad50\ubd80 \u2014 \uc591\uc790 \ud611\uc815 \uacbd\ub85c",
        ], PURPLE),
    ]
    y_offset = Cm(4)
    for zone_name, items, color in zones:
        add_rounded_rect(slide, Cm(1.5), y_offset, Cm(24), Cm(0.8 + len(items) * 0.65), SURFACE, border_color=color)
        add_textbox(slide, Cm(2), y_offset + Cm(0.1), Cm(4), Cm(0.5),
                    zone_name, font_size=9, font_color=color, bold=True)
        for j, item in enumerate(items):
            add_textbox(slide, Cm(5), y_offset + Cm(0.1 + j * 0.65), Cm(19), Cm(0.5),
                        "\u2022 " + item, font_size=8, font_color=BLACK)
        y_offset += Cm(1.2 + len(items) * 0.65)

    # ================================================================
    # SLIDE 13: 도입 H/W 사양
    # ================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_LIGHT)
    add_page_header(slide, "\ub3c4\uc785 H/W \uc0ac\uc591 \ubc0f \uc218\ub7c9", "Part 2 \u2014 \uc7a5\ube44 \uad6c\uc131")

    hw_rows = [
        ("\uc911\uacc4\uc11c\ubc84", "CPU 3.2GHz 16Core \xb7 Memory 32GB \xb7 DISK 960GB SSD \xd72EA \xb7 NIC 10GbE 2Port \xb7 \uc774\uc911\ud654 \xb7 Windows", "2\uc2dd", "DMZ \xb7 Active-Standby"),
        ("PKD\uc11c\ubc84", "CPU 3.2GHz 16Core \xb7 Memory 32GB \xb7 DISK 960GB SSD \xd72EA \xb7 NIC 10GbE 2Port \xb7 \uc774\uc911\ud654 \xb7 Windows", "2\uc2dd", "\uc11c\ube44\uc2a4 \xb7 Active-Active"),
        ("\uad00\ub9ac\uc11c\ubc84", "CPU 3.2GHz 16Core \xb7 Memory 32GB \xb7 DISK 960GB SSD \xd72EA \xb7 NIC 10GbE 2Port \xb7 Windows", "1\uc2dd", "API Gateway \xb7 \uad00\ub9ac UI"),
        ("\uc2a4\ud1a0\ub9ac\uc9c0", "NVMe SSD 7.68TB \xd728EA \xb7 Usable 150.42TiB \xb7 RaidB \xb7 32G 4p HBA \xd74ea", "1\uc2dd", "SAN \uacf5\uc720 \uc2a4\ud1a0\ub9ac\uc9c0"),
        ("SAN \uc2a4\uc704\uce58", "24port 32Gb FC Switch", "2\uc2dd", "Fabric A/B \uc774\uc911\ud654"),
        ("\ub799 \ubc0f KVM", "Deep Rack \xb7 16port KVM \uc9c0\uc6d0", "2\uc2dd", ""),
    ]
    tbl = add_table(slide, Cm(1.5), Cm(4), Cm(24), Cm(6.5), len(hw_rows) + 1, 4)
    tbl.columns[0].width = Cm(3.5)
    tbl.columns[1].width = Cm(13.5)
    tbl.columns[2].width = Cm(2)
    tbl.columns[3].width = Cm(5)
    set_cell(tbl, 0, 0, "\ud488\uba85", bold=True, fill_color=NAVY, font_color=WHITE, font_size=8, alignment=PP_ALIGN.CENTER)
    set_cell(tbl, 0, 1, "\uc138\ubd80 \uc0ac\uc591", bold=True, fill_color=NAVY, font_color=WHITE, font_size=8, alignment=PP_ALIGN.CENTER)
    set_cell(tbl, 0, 2, "\uc218\ub7c9", bold=True, fill_color=NAVY, font_color=WHITE, font_size=8, alignment=PP_ALIGN.CENTER)
    set_cell(tbl, 0, 3, "\ube44\uace0", bold=True, fill_color=NAVY, font_color=WHITE, font_size=8, alignment=PP_ALIGN.CENTER)
    for i, (name, spec, qty, note) in enumerate(hw_rows):
        r = i + 1
        bg = ALT_ROW if r % 2 == 1 else None
        set_cell(tbl, r, 0, name, bold=True, font_size=8, fill_color=bg, alignment=PP_ALIGN.CENTER)
        set_cell(tbl, r, 1, spec, font_size=7, fill_color=bg)
        set_cell(tbl, r, 2, qty, bold=True, font_size=8, fill_color=bg, alignment=PP_ALIGN.CENTER)
        set_cell(tbl, r, 3, note, font_size=7, fill_color=bg)

    # ================================================================
    # SLIDE 14: 도입 S/W 사양
    # ================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_LIGHT)
    add_page_header(slide, "\ub3c4\uc785 S/W \uc0ac\uc591 \ubc0f \uc218\ub7c9", "Part 2 \u2014 \uc7a5\ube44 \uad6c\uc131")

    sw_rows = [
        (True, "\ud1b5\ud569 \ubc1c\uc8fc", "", "", ""),
        (False, "PKD \uc911\uacc4 \uc5d0\uc774\uc804\ud2b8", "ICAO PKD \uc11c\ubc84 \uc5f0\uacc4 \xb7 DS\uc778\uc99d\uc11c/CRL \uc790\ub3d9 \ub2e4\uc6b4\ub85c\ub4dc \xb7 CSR \uae30\ubc18 \uc778\uc99d \xb7 LDAP V3 + REST", "2\uc2dd", "\uc911\uacc4\uc11c\ubc84"),
        (False, "PKD Manager", "PKI \uae30\ubc18 DS\uc778\uc99d\uc11c/CRL \uad00\ub9ac \xb7 ICAO \uaddc\uc815 PA \uc778\uc99d \xb7 \uc0ac\uc6a9\uc790/\uc778\uc99d\uc11c \uad00\ub9ac \xb7 \ud1b5\uacc4 \ubc0f \uac10\uc2dc", "1\uc2dd", "PKD\uc11c\ubc84"),
        (False, "ICAO PKI Toolkit", "\uc804\uc790\uc5ec\uad8c \ud310\ub3c5 \uae30\ub2a5 (\ub2e8\ub9d0\ubcc4 \ub77c\uc774\uc120\uc2a4 1,500\ub300) \xb7 PA/BAC/AA \uc778\uc99d", "1\uc2dd", "PKD\uc11c\ubc84 + \uc2ec\uc0ac \ub2e8\ub9d0"),
        (False, "LDAP Server", "CSCA, DS, CRL \uac8c\uc2dc \xb7 LDAP v2/v3 \xb7 UTF-8 (3,000 Entry)", "2\uc2dd", "PKD\uc11c\ubc84 #1, #2"),
        (False, "DBMS (Oracle)", "8 Core (Named 200 Users) \xb7 Stored Procedure \xb7 XML", "1\uc2dd", "\uc2a4\ud1a0\ub9ac\uc9c0 \uc5f0\uacb0"),
        (True, "\uc0c1\uc6a9 S/W \uc9c1\uc811 \uad6c\ub9e4", "", "", ""),
        (False, "\uc11c\ubc84\ubcf4\uc548", "Secuve TOS V5.0 for Windows", "5\uc2dd", "\uc804\uccb4 \uc11c\ubc84"),
        (False, "\ubaa8\ub2c8\ud130\ub9c1", "Sycros v3 Agent & System Monitoring", "5\uc2dd", "\uc804\uccb4 \uc11c\ubc84"),
        (False, "DB \uc811\uadfc\uc81c\uc5b4", "Chakra MAX \u2014 7~8 Core", "1\uc2dd", "DB\uc11c\ubc84"),
    ]
    num_data = len(sw_rows) + 1
    tbl = add_table(slide, Cm(1.5), Cm(4), Cm(24), Cm(8), num_data, 4)
    tbl.columns[0].width = Cm(4.5)
    tbl.columns[1].width = Cm(12.5)
    tbl.columns[2].width = Cm(2)
    tbl.columns[3].width = Cm(5)
    set_cell(tbl, 0, 0, "\ud488\uba85", bold=True, fill_color=NAVY, font_color=WHITE, font_size=8, alignment=PP_ALIGN.CENTER)
    set_cell(tbl, 0, 1, "\uaddc\uaca9 \ubc0f \uae30\ub2a5", bold=True, fill_color=NAVY, font_color=WHITE, font_size=8, alignment=PP_ALIGN.CENTER)
    set_cell(tbl, 0, 2, "\uc218\ub7c9", bold=True, fill_color=NAVY, font_color=WHITE, font_size=8, alignment=PP_ALIGN.CENTER)
    set_cell(tbl, 0, 3, "\uc124\uce58 \uc704\uce58", bold=True, fill_color=NAVY, font_color=WHITE, font_size=8, alignment=PP_ALIGN.CENTER)
    for i, row in enumerate(sw_rows):
        r = i + 1
        if row[0]:  # section header
            set_cell(tbl, r, 0, row[1], bold=True, font_size=7, font_color=GRAY, fill_color=RGBColor(0xF1, 0xF5, 0xF9))
            set_cell(tbl, r, 1, "", fill_color=RGBColor(0xF1, 0xF5, 0xF9))
            set_cell(tbl, r, 2, "", fill_color=RGBColor(0xF1, 0xF5, 0xF9))
            set_cell(tbl, r, 3, "", fill_color=RGBColor(0xF1, 0xF5, 0xF9))
            merge_cells(tbl, r, 0, r, 3)
        else:
            bg = ALT_ROW if r % 2 == 0 else None
            set_cell(tbl, r, 0, row[1], bold=True, font_size=7, fill_color=bg, alignment=PP_ALIGN.CENTER)
            set_cell(tbl, r, 1, row[2], font_size=7, fill_color=bg)
            set_cell(tbl, r, 2, row[3], bold=True, font_size=8, fill_color=bg, alignment=PP_ALIGN.CENTER)
            set_cell(tbl, r, 3, row[4], font_size=7, fill_color=bg)

    # ================================================================
    # SLIDE 15: Local PKD 구축 전략
    # ================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_LIGHT)
    add_page_header(slide, "Local PKD \uad6c\ucd95 \uc804\ub7b5 \u2014 \uc778\uc99d\uc11c \uc218\uc9d1 \uacbd\ub85c", "Part 2 \u2014 \ud575\uc2ec \uae30\ub2a5")

    # Left: 수집 경로 table
    coll_rows = [
        ("\uc678\uad50\uc801 \uacbd\ub85c (\uc591\uc790 \ud611\uc815)", "CSCA / DSC / JRL", "\uacf5\ubb38, \uc800\uc7a5\ub9e4\uccb4 \uc218\uc2e0"),
        ("ICAO PKD (2026.03~ \uc2e0\uaddc)", "DSC / CRL / ML / DL / DSC_NC", "CSR \uc778\uc99d \ub85c\uadf8\uc778 + LDAP/REST API"),
        ("ICAO Web", "ML (\ud604\uc7ac \ud615\ud0dc\ub85c \uacf5\uac1c)", "\ud30c\uc77c \ub2e4\uc6b4\ub85c\ub4dc"),
    ]
    tbl = add_table(slide, Cm(1.5), Cm(4), Cm(12.5), Cm(3.5), len(coll_rows) + 1, 3)
    tbl.columns[0].width = Cm(4.5)
    tbl.columns[1].width = Cm(4)
    tbl.columns[2].width = Cm(4)
    set_cell(tbl, 0, 0, "\uc218\uc9d1 \uacbd\ub85c", bold=True, fill_color=NAVY, font_color=WHITE, font_size=8, alignment=PP_ALIGN.CENTER)
    set_cell(tbl, 0, 1, "\uc815\ubcf4 \uc885\ub958", bold=True, fill_color=NAVY, font_color=WHITE, font_size=8, alignment=PP_ALIGN.CENTER)
    set_cell(tbl, 0, 2, "\uc218\uc9d1 \ud615\ud0dc", bold=True, fill_color=NAVY, font_color=WHITE, font_size=8, alignment=PP_ALIGN.CENTER)
    for i, (src, info, method) in enumerate(coll_rows):
        r = i + 1
        bg = ALT_ROW if r % 2 == 1 else None
        set_cell(tbl, r, 0, src, bold=True, font_size=7, fill_color=bg)
        set_cell(tbl, r, 1, info, font_size=7, fill_color=bg)
        set_cell(tbl, r, 2, method, font_size=7, fill_color=bg)

    # Right: flow + stats + formats
    add_textbox(slide, Cm(15), Cm(4), Cm(10), Cm(0.8),
                "\uc678\uad50\uc801 \uacbd\ub85c + ICAO PKD + Web \u2192 Local PKD",
                font_size=10, font_color=BLUE, bold=True, alignment=PP_ALIGN.CENTER)

    add_stat_card(slide, Cm(15), Cm(5.2), Cm(5), Cm(2), "31,212", "\ucd1d \uc778\uc99d\uc11c", NAVY)
    add_stat_card(slide, Cm(20.5), Cm(5.2), Cm(5), Cm(2), "139", "\uad6d\uac00", GREEN)

    add_textbox(slide, Cm(15), Cm(7.8), Cm(10), Cm(0.5),
                "\uc9c0\uc6d0 \ud30c\uc77c \ud615\uc2dd", font_size=10, font_color=BLACK, bold=True)
    formats = "LDIF  \xb7  ML (.ml)  \xb7  PEM  \xb7  DER  \xb7  P7B  \xb7  CER  \xb7  CRL  \xb7  DL  \xb7  JRL (.bin)"
    add_textbox(slide, Cm(15), Cm(8.4), Cm(10), Cm(0.8),
                formats, font_size=9, font_color=BLUE)

    # Bottom: 인증서 분포 cards
    add_textbox(slide, Cm(1.5), Cm(8.5), Cm(24), Cm(0.5),
                "\uc6b4\uc601 \ub370\uc774\ud130 \ud604\ud669 \u2014 139\uac1c\uad6d, 31,212\uac74 (DB-LDAP 100% \ub3d9\uae30\ud654)",
                font_size=10, font_color=NAVY, bold=True)
    cert_stats = [
        ("DSC", "29,838\uac74", "95.6%"),
        ("CSCA", "845\uac74", "2.7%"),
        ("DSC_NC", "502\uac74", "1.6%"),
        ("MLSC", "27\uac74", "0.1%"),
        ("CRL", "69\uac74", ""),
    ]
    for i, (tp, cnt, pct) in enumerate(cert_stats):
        x = Cm(1.5 + i * 5)
        label = f"{tp} ({pct})" if pct else tp
        add_stat_card(slide, x, Cm(9.3), Cm(4.5), Cm(2.5), cnt, label, NAVY)

    # ================================================================
    # SLIDE 16: 핵심 기능 요약
    # ================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_LIGHT)
    add_page_header(slide, "\ud575\uc2ec \uae30\ub2a5 \uc694\uc57d", "Part 2 \u2014 \uae30\ub2a5 \uc124\uacc4")

    feat_rows = [
        ("\uc778\uc99d\uc11c \ud1b5\ud569 \uad00\ub9ac", "LDIF/ML/PEM/DER/P7B/DL/CRL \uc6d0\ud074\ub9ad \uc5c5\ub85c\ub4dc + \uc790\ub3d9 \ud30c\uc2f1\xb7\uac80\uc99d\xb7\uc800\uc7a5", "\uc6b4\uc601 90% \ub2e8\ucd95"),
        ("PA \uac80\uc99d", "ICAO 9303 Part 10 & 11 \uc644\ubcbd \uc900\uc218 8\ub2e8\uacc4 \uc790\ub3d9 \uc704\xb7\ubcc0\uc870 \uac80\uc99d", "99.9% \uc815\ud655\ub3c4"),
        ("\uc778\uc99d\uc11c \uac80\uc0c9", "139\uac1c\uad6d 3\ub9cc+ \uc778\uc99d\uc11c \ucd08\uace0\uc18d \uac80\uc0c9 (\uad6d\uac00/\ud0c0\uc785/\uc0c1\ud0dc/DN \ud544\ud130)", "40ms \uc751\ub2f5"),
        ("DB-LDAP \ub3d9\uae30\ud654", "30\ucd08 \uc8fc\uae30 \uc790\ub3d9 \uccb4\ud06c, \ubd88\uc77c\uce58 \uc989\uc2dc \uac10\uc9c0 \ubc0f \uc790\ub3d9 \uc870\uc815", "100% \uc77c\uad00\uc131"),
        ("ICAO \ubc84\uc804 \uad00\ub9ac", "\ub9e4\uc77c \uc790\ub3d9 \uccb4\ud06c, \uc2e0\uaddc \ubc84\uc804 \uac10\uc9c0 \uc2dc \uc2e4\uc2dc\uac04 \uc54c\ub9bc", "\uc790\ub3d9\ud654"),
        ("CSR \uad00\ub9ac", "ICAO PKD CSR \uc0dd\uc131 (RSA-2048, SHA256withRSA, PKCS#10)", "PKD \uc804\ud658 \ub300\uc751"),
        ("DL\xb7Non-Conformant \uad00\ub9ac", "\uac01\uad6d DL \ubc0f \ube44\ud45c\uc900 DSC \ubbf8\uc900\uc218 \uc0ac\ud56d \uc870\ud68c\xb7\uad00\ub9ac, DSC_NC \ubcf4\uace0\uc11c", "\uc0ac\uc5c5\uacfc\uc81c #11"),
        ("ICRM \uc5f0\ub3d9 API", "REST API + X-API-Key \uc778\uc99d, Export (DER/PEM/ZIP), \uc2e4\uc2dc\uac04 PA", "\uc0ac\uc5c5\uacfc\uc81c #15"),
    ]
    tbl = add_table(slide, Cm(1.5), Cm(3.8), Cm(24), Cm(7.5), len(feat_rows) + 1, 3)
    tbl.columns[0].width = Cm(5)
    tbl.columns[1].width = Cm(13)
    tbl.columns[2].width = Cm(6)
    set_cell(tbl, 0, 0, "\uae30\ub2a5", bold=True, fill_color=NAVY, font_color=WHITE, font_size=8, alignment=PP_ALIGN.CENTER)
    set_cell(tbl, 0, 1, "\uc124\uba85", bold=True, fill_color=NAVY, font_color=WHITE, font_size=8, alignment=PP_ALIGN.CENTER)
    set_cell(tbl, 0, 2, "\ud6a8\uacfc", bold=True, fill_color=NAVY, font_color=WHITE, font_size=8, alignment=PP_ALIGN.CENTER)
    for i, (feat, desc, effect) in enumerate(feat_rows):
        r = i + 1
        bg = ALT_ROW if r % 2 == 1 else None
        is_highlight = (r == 2)  # PA 검증
        set_cell(tbl, r, 0, feat, bold=True, font_size=8, fill_color=bg)
        set_cell(tbl, r, 1, desc, font_size=7, fill_color=bg)
        set_cell(tbl, r, 2, effect, font_size=7, font_color=GREEN, bold=True, fill_color=bg, alignment=PP_ALIGN.CENTER)

    # ================================================================
    # SLIDE 17: 운영 실적 (인증서 분포)
    # ================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_LIGHT)
    add_page_header(slide, "ICAO PKD Download \uc0ac\uc774\ud2b8 \uac8c\uc2dc \ub370\uc774\ud130 \ub0b4\uc758 \uc778\uc99d\uc11c \ubd84\ud3ec \ud604\ud669", "Part 2 \u2014 \uac80\uc99d\ub41c \uc131\uacfc")

    # Certificate type cards
    cert_types = [
        ("DSC 29,838\uac74", "95.6% \u2014 \ubb38\uc11c \uc11c\uba85 \uc778\uc99d\uc11c (Document Signer Certificate)", BLUE),
        ("CSCA 845\uac74", "2.7% \u2014 \uad6d\uac00 \uc11c\uba85 \uc778\uc99d \uae30\uad00 (Country Signing CA)", GREEN),
        ("DSC_NC 502\uac74", "1.6% \u2014 \ube44\ud45c\uc900 \uc778\uc99d\uc11c (Non-Conformant DSC)", ORANGE),
        ("MLSC 27\uac74", "0.1% \u2014 \ub9c8\uc2a4\ud130\ub9ac\uc2a4\ud2b8 \uc11c\uba85 \uc778\uc99d\uc11c (ML Signer Certificate)", PURPLE),
        ("CRL 69\uac74", "\uc778\uc99d\uc11c \ud3d0\uae30 \ubaa9\ub85d (Certificate Revocation List)", GREEN),
    ]
    for i, (title, desc, color) in enumerate(cert_types):
        y = Cm(4.2 + i * 1.8)
        add_rect(slide, Cm(1.5), y, Cm(0.12), Cm(1.4), color)
        add_textbox(slide, Cm(2.2), y + Cm(0.1), Cm(23), Cm(0.5),
                    title, font_size=11, font_color=BLACK, bold=True)
        add_textbox(slide, Cm(2.2), y + Cm(0.7), Cm(23), Cm(0.5),
                    desc, font_size=9, font_color=GRAY)

    # Total card
    add_stat_card(slide, Cm(1.5), Cm(13.2), Cm(6), Cm(2), "31,212", "\ucd1d \uc778\uc99d\uc11c (139\uac1c\uad6d)", NAVY)
    add_stat_card(slide, Cm(8), Cm(13.2), Cm(6), Cm(2), "100%", "DB-LDAP \ub3d9\uae30\ud654\uc728", GREEN)

    # ================================================================
    # SLIDE 18: 성능 벤치마크
    # ================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_LIGHT)
    add_page_header(slide, "\uc131\ub2a5 \ubca4\uce58\ub9c8\ud06c", "Part 2 \u2014 \uc131\ub2a5")

    bench_data = [
        ("PA \uac80\uc99d (8\ub2e8\uacc4)", "100~300ms", "ICAO 9303 \uc644\ubcbd \uc900\uc218"),
        ("\uc778\uc99d\uc11c \uac80\uc0c9 (3\ub9cc\uac74)", "40ms", "1,975\ubc30 \uac1c\uc120"),
        ("Trust Chain \uac80\uc99d", "20ms", "CSCA \u2192 LC \u2192 DSC \ub2e4\ub2e8\uacc4"),
        ("DB\u2194LDAP \ub3d9\uae30\ud654 \ud655\uc778", "1~2\ucd08", "31,000+ \uc778\uc99d\uc11c \ube44\uad50"),
        ("Master List \uc5c5\ub85c\ub4dc (537\uac74)", "3\ubd84 39\ucd08", "\ud30c\uc2f1+\uac80\uc99d+DB+LDAP"),
        ("\uc804\uccb4 Export (ZIP)", "2\ucd08", "DIT \uad6c\uc870 \ud3f4\ub354"),
    ]
    tbl = add_table(slide, Cm(1.5), Cm(4), Cm(24), Cm(5.5), len(bench_data) + 1, 3)
    tbl.columns[0].width = Cm(8)
    tbl.columns[1].width = Cm(6)
    tbl.columns[2].width = Cm(10)
    set_cell(tbl, 0, 0, "\uc791\uc5c5", bold=True, fill_color=NAVY, font_color=WHITE, font_size=9, alignment=PP_ALIGN.CENTER)
    set_cell(tbl, 0, 1, "\ucc98\ub9ac \uc2dc\uac04", bold=True, fill_color=NAVY, font_color=WHITE, font_size=9, alignment=PP_ALIGN.CENTER)
    set_cell(tbl, 0, 2, "\ube44\uace0", bold=True, fill_color=NAVY, font_color=WHITE, font_size=9, alignment=PP_ALIGN.CENTER)
    for i, (task, time_val, note) in enumerate(bench_data):
        r = i + 1
        bg = ALT_ROW if r % 2 == 1 else None
        set_cell(tbl, r, 0, task, bold=True, font_size=9, fill_color=bg)
        set_cell(tbl, r, 1, time_val, font_size=9, font_color=GREEN, bold=True, fill_color=bg, alignment=PP_ALIGN.CENTER)
        set_cell(tbl, r, 2, note, font_size=9, fill_color=bg)

    # Bottom stat cards
    add_stat_card(slide, Cm(5), Cm(10.5), Cm(7), Cm(2.5), "<200ms", "\ud3c9\uade0 API \uc751\ub2f5 \uc2dc\uac04", GREEN)
    add_stat_card(slide, Cm(14), Cm(10.5), Cm(7), Cm(2.5), "99.9%+", "\uc2dc\uc2a4\ud15c \uac00\uc6a9\uc131 (SLA)", GREEN)

    # ================================================================
    # SLIDE 19: Part 3 간지
    # ================================================================
    add_part_divider(prs, blank_layout, 3,
                     "\ubcf4\uc548 \ub300\ucc45 \ubc0f \ud504\ub85c\uc81d\ud2b8 \uad00\ub9ac",
                     "\ubc95\uc801 \ucef4\ud50c\ub77c\uc774\uc5b8\uc2a4 \xb7 \ubcf4\uc548 \uc778\uc99d \xb7 \uccb4\uacc4\uc801 \uc0ac\uc5c5 \uad00\ub9ac")

    # ================================================================
    # SLIDE 20: 보안 대책 및 법적 준수
    # ================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_LIGHT)
    add_page_header(slide, "\ubcf4\uc548 \ub300\ucc45 \ubc0f \ubc95\uc801 \uc900\uc218", "Part 3 \u2014 \ubcf4\uc548")

    # Left: 기술적 보안
    add_textbox(slide, Cm(1.5), Cm(3.8), Cm(12), Cm(0.6),
                "\uae30\uc220\uc801 \ubcf4\uc548", font_size=12, font_color=BLACK, bold=True)
    sec_items = [
        ("[1] AES-256-GCM \uac1c\uc778\uc815\ubcf4 \uc554\ud638\ud654", "\uc5ec\uad8c\ubc88\ud638\xb7\uc131\uba85\xb7\uc774\uba54\uc77c \ub4f1 PII \ud544\ub4dc \uc778\uc99d\ub41c \uc554\ud638\ud654"),
        ("[2] JWT + RBAC + API Key \uc778\uc99d", "3\uc911 \uc778\uc99d \uccb4\uacc4, Role \uae30\ubc18 \uc811\uadfc \uc81c\uc5b4"),
        ("[3] \uc774\uc911 \uac10\uc0ac \ub85c\uadf8", "\uc778\uc99d \uac10\uc0ac(auth_audit) + \uc6b4\uc601 \uac10\uc0ac(operation_audit)"),
        ("[4] 100% \ud30c\ub77c\ubbf8\ud130\ud654 SQL", "SQL \uc778\uc81d\uc158 \uc6d0\ucc9c \ucc28\ub2e8, CSP \ud5e4\ub354 XSS \ubc29\uc5b4"),
    ]
    for i, (title, desc) in enumerate(sec_items):
        y = Cm(4.6 + i * 1.5)
        add_textbox(slide, Cm(1.5), y, Cm(12), Cm(0.5),
                    title, font_size=9, font_color=BLACK, bold=True)
        add_textbox(slide, Cm(1.5), y + Cm(0.5), Cm(12), Cm(0.6),
                    desc, font_size=8, font_color=GRAY)

    # Right: 법적 컴플라이언스 table
    add_textbox(slide, Cm(14), Cm(3.8), Cm(12), Cm(0.6),
                "\ubc95\uc801 \ucef4\ud50c\ub77c\uc774\uc5b8\uc2a4", font_size=12, font_color=BLACK, bold=True)
    law_rows = [
        ("\uac1c\uc778\uc815\ubcf4\ubcf4\ud638\ubc95 \xa729 \uc548\uc804\uc870\uce58", "\uc900\uc218 \u2014 AES-256-GCM \uc554\ud638\ud654"),
        ("\uac1c\uc778\uc815\ubcf4\ubcf4\ud638\ubc95 \xa724 \uace0\uc720\uc2dd\ubcc4\uc815\ubcf4", "\uc900\uc218 \u2014 \uc5ec\uad8c\ubc88\ud638 \uc554\ud638\ud654 \uc800\uc7a5"),
        ("\uc804\uc790\uc815\ubd80\ubc95 \uc2dc\ud589\ub839 \xa769 KCMVP", "\ub300\uc751 \uac00\ub2a5 \u2014 OpenSSL FIPS \ubaa8\ub4dc"),
        ("CC \uc778\uc99d (\ubcf4\uc548\uae30\ub2a5\ud655\uc778\uc11c)", "\ub300\uc751 \uac00\ub2a5 \u2014 \ud3c9\uac00 \ub300\ube44 \uc124\uacc4"),
        ("ICAO Doc 9303 Part 10 & 11", "\uc644\ubcbd \uc900\uc218"),
    ]
    tbl = add_table(slide, Cm(14), Cm(4.6), Cm(11.5), Cm(5), len(law_rows) + 1, 2)
    tbl.columns[0].width = Cm(5.5)
    tbl.columns[1].width = Cm(6)
    set_cell(tbl, 0, 0, "\ubc95\ub839/\ud45c\uc900", bold=True, fill_color=NAVY, font_color=WHITE, font_size=8, alignment=PP_ALIGN.CENTER)
    set_cell(tbl, 0, 1, "\uc900\uc218 \uc0c1\ud0dc", bold=True, fill_color=NAVY, font_color=WHITE, font_size=8, alignment=PP_ALIGN.CENTER)
    for i, (law, status) in enumerate(law_rows):
        r = i + 1
        bg = ALT_ROW if r % 2 == 1 else None
        s_color = GREEN if "\uc900\uc218" in status else BLUE
        set_cell(tbl, r, 0, law, bold=True, font_size=7, fill_color=bg)
        set_cell(tbl, r, 1, status, font_size=7, font_color=s_color, bold=True, fill_color=bg, alignment=PP_ALIGN.CENTER)

    # ================================================================
    # SLIDE 21: 단계별 추진 일정
    # ================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_LIGHT)
    add_page_header(slide, "\ub2e8\uacc4\ubcc4 \ucd94\uc9c4 \uc77c\uc815", "Part 3 \u2014 \ud504\ub85c\uc81d\ud2b8 \uad00\ub9ac")

    phases = [
        ("1", "\ud658\uacbd \uad6c\ucd95 \ubc0f \uc124\uce58", "1~2\uc8fc",
         "Docker/Podman \ucee8\ud14c\uc774\ub108 \ubc30\ud3ec \xb7 DB \uc2a4\ud0a4\ub9c8 \ucd08\uae30\ud654 \xb7 LDAP DIT \uad6c\uc131 \xb7 SSL \uc778\uc99d\uc11c \uc124\uc815", BLUE),
        ("2", "\ub370\uc774\ud130 \uad6c\ucd95 \ubc0f \uac80\uc99d", "2~4\uc8fc",
         "ICAO PKD \ub2e4\uc6b4\ub85c\ub4dc \xb7 LDIF/Master List \uc5c5\ub85c\ub4dc \xb7 Trust Chain \uac80\uc99d \xb7 DB\u2194LDAP \ub3d9\uae30\ud654 \ud655\uc778", ORANGE),
        ("3", "\uc5f0\ub3d9 \ubc0f \uc2dc\ubc94 \uc6b4\uc601", "2~4\uc8fc",
         "\ucd9c\uc785\uad6d \uc2ec\uc0ac \uc2dc\uc2a4\ud15c API \uc5f0\ub3d9 \xb7 PA \uac80\uc99d \uc5f0\ub3d9 \ud14c\uc2a4\ud2b8 \xb7 \uad00\ub9ac\uc790 \uad50\uc721 \xb7 \uc2dc\ubc94 \uc6b4\uc601", GREEN),
        ("4", "\ubcf8\uaca9 \uc6b4\uc601 \ubc0f \uace0\ub3c4\ud654", "\uc0c1\uc2dc",
         "ICAO PKD \uc815\uc2dd \ud68c\uc6d0 \ub4f1\ub85d \xb7 \uc591\uc790 \ud611\uc815 \ucd94\uc9c4 \xb7 \uc11c\ube44\uc2a4 \ubc94\uc704 \ub2e8\uacc4\uc801 \ud655\uc7a5", PURPLE),
    ]
    for i, (num, title, period, desc, color) in enumerate(phases):
        y = Cm(4 + i * 2.5)
        add_rounded_rect(slide, Cm(1.5), y, Cm(1.2), Cm(1.2), color,
                         num, font_size=14, bold=True)
        add_textbox(slide, Cm(3.2), y + Cm(0.05), Cm(12), Cm(0.5),
                    title, font_size=11, font_color=BLACK, bold=True)
        add_rounded_rect(slide, Cm(13), y + Cm(0.05), Cm(2.5), Cm(0.6), color,
                         period, font_size=7, bold=True)
        add_textbox(slide, Cm(3.2), y + Cm(0.7), Cm(20), Cm(0.8),
                    desc, font_size=8, font_color=GRAY, line_spacing=1.4)

    # Total
    add_rounded_rect(slide, Cm(5), Cm(14.2), Cm(16), Cm(0.8), SURFACE, border_color=BORDER)
    add_textbox(slide, Cm(5.5), Cm(14.3), Cm(15), Cm(0.6),
                "Phase 1~3 \ucd1d \uc18c\uc694 \uae30\uac04: 6~10\uc8fc (\ubc1c\uc8fc\ucc98 \ud658\uacbd \ubc0f \uc694\uad6c\uc0ac\ud56d\uc5d0 \ub530\ub77c \uc870\uc815)",
                font_size=10, font_color=BLUE, bold=True, alignment=PP_ALIGN.CENTER)

    # ================================================================
    # SLIDE 22: 유지보수 및 운영 지원
    # ================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_LIGHT)
    add_page_header(slide, "\uc720\uc9c0\ubcf4\uc218 \ubc0f \uc6b4\uc601 \uc9c0\uc6d0 \uccb4\uacc4", "Part 3 \u2014 \uc6b4\uc601 \uc9c0\uc6d0")

    # 3 service cards
    maint_cards = [
        ("\uc0c1\uc2dc \ubaa8\ub2c8\ud130\ub9c1", "\uc2dc\uc2a4\ud15c \uba54\ud2b8\ub9ad \uc2e4\uc2dc\uac04 \uc218\uc9d1\n\uc11c\ube44\uc2a4 \ud5ec\uc2a4\uccb4\ud06c \uc790\ub3d9\ud654\nSSE \uae30\ubc18 \uc2e4\uc2dc\uac04 \uc54c\ub9bc\nDB\u2194LDAP \ub3d9\uae30\ud654 \ub300\uc2dc\ubcf4\ub4dc", BLUE),
        ("\uc7a5\uc560 \ub300\uc751", "\uc7a5\uc560 \ub4f1\uae09\ubcc4 \ub300\uc751 \uc808\ucc28 (4\ub2e8\uacc4)\n\uc790\ub3d9 \uc7ac\uc2dc\ub3c4 \ubc0f \ubcf5\uad6c \uba54\ucee4\ub2c8\uc998\nLDAP \uc7a5\uc560 \uc2dc DB-only \ubaa8\ub4dc\nReconciliation \uc790\ub3d9 \ubcf5\uad6c", ORANGE),
        ("\uc6b4\uc601 \ubcf4\uace0", "\uc778\uc99d\uc11c \ud604\ud669 \ubcf4\uace0\uc11c \uc790\ub3d9 \uc0dd\uc131\n\uac80\uc99d \uacb0\uacfc \ud1b5\uacc4 \ubcf4\uace0\uc11c\n\uac10\uc0ac \ub85c\uadf8 \ud1b5\uacc4 \uc81c\uacf5\n\uc6d4\uac04/\ubd84\uae30 \uc6b4\uc601 \ubcf4\uace0\uc11c", GREEN),
    ]
    for i, (title, desc, color) in enumerate(maint_cards):
        x = Cm(1.5 + i * 8.5)
        add_rect(slide, x, Cm(4), Cm(7.8), Cm(0.12), color)
        add_rounded_rect(slide, x, Cm(4.12), Cm(7.8), Cm(3.8), SURFACE, border_color=BORDER)
        add_textbox(slide, x + Cm(0.5), Cm(4.5), Cm(6.8), Cm(0.5),
                    title, font_size=10, font_color=BLACK, bold=True, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x + Cm(0.5), Cm(5.2), Cm(6.8), Cm(2.5),
                    desc, font_size=8, font_color=GRAY, line_spacing=1.5, alignment=PP_ALIGN.CENTER)

    # SLA table
    add_textbox(slide, Cm(1.5), Cm(8.5), Cm(12), Cm(0.5),
                "SLA \uc218\uc900", font_size=10, font_color=BLACK, bold=True)
    sla_rows = [
        ("\uc2dc\uc2a4\ud15c \uac00\uc6a9\uc131", "99.9% \uc774\uc0c1"),
        ("PA \uac80\uc99d \uc751\ub2f5 \uc2dc\uac04", "300ms \uc774\ub0b4"),
        ("\uc7a5\uc560 \ucd08\uae30 \ub300\uc751", "30\ubd84 \uc774\ub0b4"),
        ("ICAO \uc778\uc99d\uc11c \uac31\uc2e0", "\ub9e4\uc77c \uc790\ub3d9 \uccb4\ud06c"),
    ]
    tbl = add_table(slide, Cm(1.5), Cm(9.1), Cm(11.5), Cm(3.5), len(sla_rows) + 1, 2)
    tbl.columns[0].width = Cm(5.5)
    tbl.columns[1].width = Cm(6)
    set_cell(tbl, 0, 0, "\ud56d\ubaa9", bold=True, fill_color=NAVY, font_color=WHITE, font_size=8, alignment=PP_ALIGN.CENTER)
    set_cell(tbl, 0, 1, "\ubaa9\ud45c", bold=True, fill_color=NAVY, font_color=WHITE, font_size=8, alignment=PP_ALIGN.CENTER)
    for i, (item, target) in enumerate(sla_rows):
        r = i + 1
        set_cell(tbl, r, 0, item, font_size=8, fill_color=ALT_ROW if r % 2 == 1 else None)
        set_cell(tbl, r, 1, target, font_size=8, bold=True, font_color=GREEN, fill_color=ALT_ROW if r % 2 == 1 else None, alignment=PP_ALIGN.CENTER)

    # Update table
    add_textbox(slide, Cm(14), Cm(8.5), Cm(12), Cm(0.5),
                "\uc5c5\ub370\uc774\ud2b8 \uccb4\uacc4", font_size=10, font_color=BLACK, bold=True)
    upd_rows = [
        ("ICAO \ud45c\uc900 \ubcc0\uacbd \ub300\uc751", "\uc989\uc2dc \ubc18\uc601"),
        ("\ubcf4\uc548 \ud328\uce58", "\ubc1c\uacac \uc989\uc2dc \uc801\uc6a9"),
        ("\uae30\ub2a5 \uc5c5\ub370\uc774\ud2b8", "\ubd84\uae30\ubcc4 \ub9b4\ub9ac\uc2a4"),
        ("DB \uc2a4\ud0a4\ub9c8 \ub9c8\uc774\uadf8\ub808\uc774\uc158", "\ubb34\uc911\ub2e8 \uc801\uc6a9"),
    ]
    tbl = add_table(slide, Cm(14), Cm(9.1), Cm(11.5), Cm(3.5), len(upd_rows) + 1, 2)
    tbl.columns[0].width = Cm(5.5)
    tbl.columns[1].width = Cm(6)
    set_cell(tbl, 0, 0, "\ud56d\ubaa9", bold=True, fill_color=NAVY, font_color=WHITE, font_size=8, alignment=PP_ALIGN.CENTER)
    set_cell(tbl, 0, 1, "\uc8fc\uae30", bold=True, fill_color=NAVY, font_color=WHITE, font_size=8, alignment=PP_ALIGN.CENTER)
    for i, (item, cycle) in enumerate(upd_rows):
        r = i + 1
        set_cell(tbl, r, 0, item, font_size=8, fill_color=ALT_ROW if r % 2 == 1 else None)
        set_cell(tbl, r, 1, cycle, font_size=8, fill_color=ALT_ROW if r % 2 == 1 else None, alignment=PP_ALIGN.CENTER)

    # ================================================================
    # SLIDE 23: 도입 기대효과
    # ================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_LIGHT)
    add_page_header(slide, "\ub3c4\uc785 \uae30\ub300\ud6a8\uacfc", "Part 3 \u2014 \uae30\ub300\ud6a8\uacfc")

    # 4 stat cards
    effect_stats = [
        ("99.9%", "\uc704\uc870 \ud0d0\uc9c0\uc728", "3\uc138\ub300 \uc815\ubc00 \uc704\uc870\uae4c\uc9c0 \uc790\ub3d9 \ud0d0\uc9c0", GREEN),
        ("90%", "\uc6b4\uc601 \ube44\uc6a9 \uc808\uac10", "\uc218\ub3d9 \uc791\uc5c5 \uc790\ub3d9\ud654, \uc778\ub825 \uc758\uc874\ub3c4 \ucd5c\uc18c\ud654", BLUE),
        ("1,200h", "\uc5f0\uac04 \uc808\uac10 \uc2dc\uac04", "\uc77c 3.5\uc2dc\uac04 \u2192 10\ubd84 \uc778\uc99d\uc11c \uad00\ub9ac", ORANGE),
        ("100%", "ICAO \uc900\uc218", "Doc 9303 \uc644\ubcbd \uc900\uc218, \uad6d\uc81c \uc694\uad6c\uc0ac\ud56d \ucda9\uc871", PURPLE),
    ]
    for i, (val, label, desc, color) in enumerate(effect_stats):
        x = Cm(1.5 + i * 6.2)
        add_rect(slide, x, Cm(4), Cm(5.7), Cm(0.1), color)
        add_rounded_rect(slide, x, Cm(4.1), Cm(5.7), Cm(2.8), SURFACE, border_color=BORDER)
        add_textbox(slide, x + Cm(0.3), Cm(4.3), Cm(5.1), Cm(0.8),
                    val, font_size=18, font_color=color, bold=True, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x + Cm(0.3), Cm(5.2), Cm(5.1), Cm(0.5),
                    label, font_size=9, font_color=BLACK, bold=True, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x + Cm(0.3), Cm(5.8), Cm(5.1), Cm(0.7),
                    desc, font_size=7, font_color=GRAY, alignment=PP_ALIGN.CENTER)

    # 정량적 효과
    add_textbox(slide, Cm(1.5), Cm(7.8), Cm(12), Cm(0.5),
                "\uc815\ub7c9\uc801 \ud6a8\uacfc", font_size=10, font_color=BLACK, bold=True)
    quant = [
        "\u2713 \ucd9c\uc785\uad6d \uc2ec\uc0ac 1\uac74\ub2f9 PA \uac80\uc99d \ucd94\uac00 \uc2dc\uac04: 0.3\ucd08 \uc774\ub0b4",
        "\u2713 \uc778\uc99d\uc11c \uc5c5\ub370\uc774\ud2b8 \uc790\ub3d9\ud654: 24\uc2dc\uac04 \ubb34\uc911\ub2e8",
        "\u2713 \uc704\uc870 \uc5ec\uad8c \uc801\ubc1c\ub960 \ud5a5\uc0c1: \ubb3c\ub9ac \uac80\uc0ac \ub300\ube44 10\ubc30+",
        "\u2713 \ud22c\uc790\uae08 \ud68c\uc218 \uae30\uac04: 6~12\uac1c\uc6d4",
    ]
    for i, item in enumerate(quant):
        add_textbox(slide, Cm(1.5), Cm(8.5 + i * 0.7), Cm(12), Cm(0.5),
                    item, font_size=8, font_color=BLACK)

    # 정성적 효과
    add_textbox(slide, Cm(14), Cm(7.8), Cm(12), Cm(0.5),
                "\uc815\uc131\uc801 \ud6a8\uacfc", font_size=10, font_color=BLACK, bold=True)
    qual = [
        "\u25c6 \uad6d\uac00 \ucd9c\uc785\uad6d \uad00\ub9ac \uad6d\uc81c \uc2e0\ub8b0\ub3c4 \ud5a5\uc0c1",
        "\u25c6 \ubd88\ubc95 \uc785\uad6d\xb7\ud14c\ub7ec \uc0ac\uc804 \ucc28\ub2e8 \uc5ed\ub7c9 \uac15\ud654",
        "\u25c6 SES \ub300\uc0c1\uad6d \ucd94\uac00 \ud655\ub300 \uae30\ubc18 \ub9c8\ub828",
        "\u25c6 \uac1c\uc778\uc815\ubcf4\ubcf4\ud638 \ubc95\uc801 \ub9ac\uc2a4\ud06c \ud574\uc18c",
    ]
    for i, item in enumerate(qual):
        add_textbox(slide, Cm(14), Cm(8.5 + i * 0.7), Cm(12), Cm(0.5),
                    item, font_size=8, font_color=BLUE)

    # ================================================================
    # SLIDE 24: 차별화 요소
    # ================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_LIGHT)
    add_page_header(slide, "\ucc28\ubcc4\ud654 \uc694\uc18c", "Part 3 \u2014 \ucc28\ubcc4\ud654")

    comp_data = [
        ("ICAO 9303 \uc900\uc218 \ubc94\uc704", "\ubd80\ubd84 \uad6c\ud604 (PA \uac80\uc99d \uc77c\ubd80 \ub204\ub77d)", "Part 10 & 11 \uc644\ubcbd \uc900\uc218 (8\ub2e8\uacc4)"),
        ("\uc778\uc99d\uc11c \uac80\uc0c9 \uc18d\ub3c4", "5~10\ucd08 \uc774\uc0c1", "40ms (1,975\ubc30 \uac1c\uc120)"),
        ("Oracle \ub124\uc774\ud2f0\ube0c \uc9c0\uc6d0", "\ub2e8\uc77c DB \uc885\uc18d \ub610\ub294 \ubbf8\uc9c0\uc6d0", "Oracle RDBMS \uc644\ubcbd \uc9c0\uc6d0"),
        ("\uc815\ubc00 \ubd84\uc11d", "\ubbf8\uc9c0\uc6d0", "\uc778\uc99d\uc11c \uc9c1\uc811 \ud30c\uc2f1\xb7\uac80\uc99d, \uad6d\uac00\ubcc4 \ud1b5\uacc4"),
        ("\uac1c\uc778\uc815\ubcf4 \uc554\ud638\ud654", "\ubbf8\uad6c\ud604 \ub610\ub294 \ubd80\ubd84", "AES-256-GCM \uc778\uc99d \uc554\ud638\ud654"),
        ("\uc544\ud0a4\ud14d\ucc98", "\ubaa8\ub178\ub9ac\uc2dd (\ud655\uc7a5 \uc81c\ud55c)", "\ub9c8\uc774\ud06c\ub85c\uc11c\ube44\uc2a4 4\uac1c \ub3c5\ub9bd \uc11c\ube44\uc2a4"),
        ("\uc790\ub3d9\ud654 \uc218\uc900", "\uc218\ub3d9 \uac1c\uc785 \ud544\uc694", "\uc644\uc804 \uc790\ub3d9 (\uc5c5\ub85c\ub4dc\xb7\uac80\uc99d\xb7\ub3d9\uae30\ud654\xb7\uc54c\ub9bc)"),
        ("\uc6b4\uc601 \ud658\uacbd", "\ud2b9\uc815 OS \uc885\uc18d", "Docker/Podman \ucee8\ud14c\uc774\ub108, RHEL 9 \uac80\uc99d"),
    ]
    tbl = add_table(slide, Cm(1.5), Cm(3.8), Cm(24), Cm(7.5), len(comp_data) + 1, 3)
    tbl.columns[0].width = Cm(5)
    tbl.columns[1].width = Cm(9)
    tbl.columns[2].width = Cm(10)
    set_cell(tbl, 0, 0, "\ud3c9\uac00 \ud56d\ubaa9", bold=True, fill_color=NAVY, font_color=WHITE, font_size=8, alignment=PP_ALIGN.CENTER)
    set_cell(tbl, 0, 1, "\uc77c\ubc18 PKD \uc194\ub8e8\uc158", bold=True, fill_color=NAVY, font_color=WHITE, font_size=8, alignment=PP_ALIGN.CENTER)
    set_cell(tbl, 0, 2, "FASTpass\u00ae PKD", bold=True, fill_color=NAVY, font_color=WHITE, font_size=8, alignment=PP_ALIGN.CENTER)
    for i, (item, gen, fast) in enumerate(comp_data):
        r = i + 1
        is_highlight = r % 2 == 1
        set_cell(tbl, r, 0, item, bold=True, font_size=8,
                 fill_color=ALT_ROW if is_highlight else None)
        set_cell(tbl, r, 1, gen, font_size=8, font_color=GRAY,
                 fill_color=ALT_ROW if is_highlight else None, alignment=PP_ALIGN.CENTER)
        set_cell(tbl, r, 2, fast, bold=True, font_size=8, font_color=BLUE,
                 fill_color=ALT_ROW if is_highlight else None, alignment=PP_ALIGN.CENTER)

    # ================================================================
    # SLIDE 25: 감사합니다 (Closing)
    # ================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, NAVY)

    add_textbox(slide, Cm(0), Cm(4.5), SLIDE_W, Cm(2),
                "\uac10\uc0ac\ud569\ub2c8\ub2e4",
                font_size=36, font_color=WHITE, bold=True,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Cm(0), Cm(7), SLIDE_W, Cm(1.2),
                "\ub300\ud55c\ubbfc\uad6d \ucd9c\uc785\uad6d \ubcf4\uc548\uc758 \uc0c8\ub85c\uc6b4 \ud45c\uc900\uc744 \uc81c\uc548\ud569\ub2c8\ub2e4",
                font_size=13, font_color=RGBColor(0x94, 0xA3, 0xB8),
                alignment=PP_ALIGN.CENTER)

    # Contact info
    add_textbox(slide, Cm(0), Cm(9.5), SLIDE_W, Cm(2),
                "\uc81c\uc548\uc0ac: (\uc8fc)\uc2a4\ub9c8\ud2b8\ucf54\uc5b4\nEmail: sales@smartcoreinc.com\nWeb: www.smartcoreinc.com",
                font_size=11, font_color=RGBColor(0x94, 0xA3, 0xB8),
                alignment=PP_ALIGN.CENTER, line_spacing=1.7)

    # Badges
    badges = "ICAO Doc 9303 \uc644\ubcbd \uc900\uc218   |   \uac1c\uc778\uc815\ubcf4\ubcf4\ud638\ubc95 \ucda9\uc871   |   KCMVP/CC \uc778\uc99d \ub300\uc751"
    add_rounded_rect(slide, Cm(4), Cm(13.5), Cm(18), Cm(0.9),
                     RGBColor(0x0A, 0x4F, 0x7A),
                     badges, font_size=8, bold=False,
                     border_color=RGBColor(0x1A, 0x6F, 0xAA))

    # Version
    add_textbox(slide, Cm(0), Cm(15.5), SLIDE_W, Cm(0.8),
                "(\uc8fc)\uc2a4\ub9c8\ud2b8\ucf54\uc5b4 \u2014 FASTpass\u00ae PKD v2.36.0",
                font_size=8, font_color=RGBColor(0x64, 0x74, 0x8B),
                alignment=PP_ALIGN.CENTER)

    # ── Save ──
    output_path = os.path.join(os.path.dirname(__file__),
                               "FASTpass_PKD_\uae30\uc220\uc81c\uc548\uc11c.pptx")
    prs.save(output_path)
    print(f"PPTX generated: {output_path}")
    print(f"   Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    build()
